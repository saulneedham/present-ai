import requests
import re
from bs4 import BeautifulSoup
import os
from urllib.parse import urljoin
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import wikipedia
from openai import OpenAI
from pptx.enum.text import PP_ALIGN
import time

#------------------------------------------------------------------------

def searchWikipedia():
    topicChosen = False

    while not topicChosen:
        topic = input("Enter a topic for your PowerPoint: ")
        try:
            results = wikipedia.search(topic, results=5)
            
            if not results:
                print(f"\nTopic not found!")
                return

            print(f"\n--- Select from the following! ---")
            
            for i, title in enumerate(results):
                print(f"[{i+1}] {title}")
            
            print("-------------------------------------------------")

            try:
                choice = int(input('Enter topic number - '))
                topicChosen = True
                topic = results[choice-1]
            except Exception as e:
                print(f"\nAn unexpected error occurred: {e}")

        except wikipedia.exceptions.DisambiguationError:
            print(f"\nThe topic '{topic}' is too general. Please be more specific.")

        except Exception as e:
            print(f"\nAn unexpected error occurred: {e}")

    # make parent folder
    parentFolder = topic.replace(' ','_')
    os.makedirs(parentFolder, exist_ok=True)
    url = "https://en.wikipedia.org/wiki/"+parentFolder

    return url, topic, parentFolder

#---------------OPENAI AGENTS----------------

def summariseCaption(captionText):
    if not captionText:
        return ""
    completion = ai.chat.completions.create(
        model='gpt-3.5-turbo',
        temperature=0.0,
        max_tokens=30,
        messages=[
            {
                'role': 'system',
                'content': (
                    "You are a caption expert. Summarize the following text into a single, "
                    "descriptive phrase between 5 and 8 words. Do not use quotes, "
                    "citations, or starting dashes. Output only the summarized phrase."
                )
            },
            {
                'role': 'user',
                'content': f"Summarize this caption: {captionText}"
            }
        ]
    )
    return completion.choices[0].message.content.strip()

def makeBulletPoints(visibleText):
    # Prepare the assistant call to produce strict JSON output:
    completion = ai.chat.completions.create(
        model='gpt-3.5-turbo',
        temperature=0.0, # deterministic
        max_tokens=400,
        messages=[
            {
                'role': 'system',
                'content': (
                    "You are an expert slide-writer that converts a chunk of text into "
                    "concise PowerPoint bullet points. STRICT RULES - follow them exactly:\n"
                    "1) Output a single joined string splitting each bullet point with a line break character.\n"
                    "2) NEVER put dashes (-) at the start of each bullet point, each bullet should be its own bit of concise info\n"
                    "3) You MUST produce between 4 and 8 bullet points. The total number of line breaks in your output CANNOT exceed 7 (for a maximum of 8 lines). This rule is non-negotiable.\n"
                    "4) The absolute maximum output length is 120 words and 600 characters.\n"
                    "5) Each bullet should be a single short sentence or phrase, ideally 6-15 words\n"
                    "6) Do not include citations, bracketed references, html, or source text.\n"
                    "7) Use plain text only; do not return markdown, lists, headings or extra fields.\n"
                    "8) If the input is short or has too little content, still return 4 concise bullets.\n"
                    "9) If you cannot identify 4 meaningful bullets, return the four best short summary phrases.\n"
                    "Tone: neutral, factual, slide-friendly.\n"
                    "Example input -> output:\n"
                    "Input: 'The Hindenburg disaster occurred in 1937 when the German passenger airship LZ 129 Hindenburg caught fire while docking in New Jersey.'\n"
                    "Output: 'Hindenburg disaster: LZ 129 caught fire while docking (1937)\nMajor loss of life and media coverage\nFire highlighted hydrogen safety risks\nInvestigation identified cause as static electricity'"
                )
            },
            {
                'role': 'user',
                'content': (
                    "Convert the following text into slide-ready bullets. "
                    "Remember the strict limits above.\n\n"
                    f"Text: {visibleText}"
                )
            }
        ]
    )
    return completion.choices[0].message.content

#---------------WEB SCRAPING FUNCTIONS----------------

def splitContent(html):
    contents = html.split('<div class="mw-heading mw-heading2"><h2 id="')
    contents.pop(0)
    referencesContent = ''

    for content in contents:
        subTopicTitle = str((content.split('"'))[0]).replace('_',' ')
        #splitting end of id tag to get subtopic title

        if subTopicTitle not in avoidedContents:
            imgs,captions = saveImages(subTopicTitle,content,avoidedImages)
            subTopicImages.append(list(zip(imgs, captions)))

            subTopicContent = re.sub(r'\[[^\]]*\]', '', removeTags(content)) #get just p tags and get rid of square brackets
            subTopicTitles.append(subTopicTitle)
            subTopicBodies.append(subTopicContent)
        elif subTopicTitle in ['References','References 2']:
            subTopicContent = extractRefs(content)
            if len(referencesContent)<2000:
                referencesContent += subTopicContent[:2000]+'\n'+'And more...'

    return referencesContent,subTopicTitles,subTopicBodies,subTopicImages

def removeTags(html):
    soup = BeautifulSoup(html, 'html.parser')
    visibleText = []

    for element in soup.find_all('p'):
        visibleText.append(element.get_text())

    return ' '.join(visibleText)


def niceRefs(html):
    maxRefs = 8

    soup = BeautifulSoup(html, 'html.parser')
    items = soup.select('ol.references > li, .reflist li')
    references = []

    for li in items:
        raw = (li.find('cite').get_text(" ", strip=True) if li.find('cite') else li.get_text(" ", strip=True))
        raw = re.sub(r'^\s*(?:\^|\^?\s*[a-z]\b(?:\s+[a-z]\b)*\s*)+', '', raw, flags=re.I)
        text = re.sub(r'\s+', ' ', raw).strip()[:300]

        url = None
        for a in li.find_all('a', href=True):
            href = a['href']
            if href.startswith('http://') or href.startswith('https://'):
                url = href
                break

        references.append((text, url))
        if len(references) == maxRefs:
            break

    return references

def extractRefs(html):
    refs = ''
    for i, (text, url) in enumerate(niceRefs(html), start=1):
        if url:
            refs+=f"{i}. {text}\n   → {url}\n"
        else:
            refs+=f"{i}. {text}\n"

    return refs

#---------------IMAGE FUNCTIONS----------------

def getImageSize(image, maxWidth=4.5, maxHeight=5.5):
    imWidth, imHeight = Image.open(image).size
    dimensions = imWidth/imHeight
    maxDimension = maxWidth/maxHeight
    
    if dimensions > maxDimension: # Landscape ratio or wider
        width = maxWidth
        height = maxWidth * (1 / dimensions)
    else: # Portrait ratio or taller
        height = maxHeight
        width = maxHeight * dimensions
        
    # Recalculate if constrained in one dimension due to the other's limit
    if width > maxWidth:
        width = maxWidth
        height = maxWidth * (1 / dimensions)
    if height > maxHeight:
        height = maxHeight
        width = maxHeight * dimensions
    
    heightLost = maxHeight - height
    widthLost = maxWidth - width
    
    return width, height, widthLost, heightLost

def saveImages(subTopicTitle,content,avoidedImages):
    soup = BeautifulSoup(content, 'html.parser')
    imagesSaved = 0
    imgs = []
    captions =[]

    topicFolder = os.path.join(parentFolder, subTopicTitle.replace(' ', '_'))
    os.makedirs(topicFolder, exist_ok=True)

    for img in soup.find_all('img'):
        if imagesSaved >= 2:
            break

        src = img.get('src') or img.get('data-src') or img.get('data-image-src') or img.get('srcset') or ''
        if not src:
            continue

        if ',' in src and ' ' in src:
            candidates = [s.strip() for s in src.split(',') if s.strip()]
            last = candidates[-1]
            src = last.split()[0]

        if src.startswith('//'):
            link = 'https:' + src
        elif src.startswith('/'):
            link = urljoin('https://en.wikipedia.org', src)
        elif src.startswith('http'):
            link = src
        else:
            continue

        if '/media/math/render/' in link: #skipping maths svgs as cant be saved as png
            continue

        # basic filename and filter check
        url_path = link.split('?', 1)[0]
        filename_only = os.path.basename(url_path)
        if any(bad in filename_only for bad in avoidedImages):
            continue

        # caption heuristics
        cap = ""
        parentFig = img.find_parent('figure')
        if parentFig:
            capTag = parentFig.find('figcaption') or parentFig.find(class_='thumbcaption')
            if capTag:
                cap = capTag.get_text(separator=' ', strip=True)
        if not cap:
            capTag = img.find_next('figcaption') or img.find_next(class_='thumbcaption')
            if capTag:
                cap = capTag.get_text(separator=' ', strip=True)
        if not cap:
            cap = img.get('alt') or img.get('title') or ""

        if cap:
            cap = summariseCaption(cap)

        ext = '.png'

        local_name = f"img{imagesSaved}{ext}"
        local_path = os.path.join(topicFolder, local_name)

        # download and save
        try:
            r = requests.get(link, headers=headers, stream=True, timeout=12)
            r.raise_for_status()
            with open(local_path, 'wb') as f:
                for chunk in r.iter_content(8192):
                    if chunk:
                        f.write(chunk)
            #print(f"Saved: {local_path}")
            imgs.append(local_path)
            captions.append(cap)
            imagesSaved += 1
        except Exception as e:
            print(f"Failed to save {link}: {e}")
            # try next image tag

    return imgs,captions

#---------------POWERPOINT FUNCTIONS----------------

def generatePP(topic,referencesContent,subTopicTitles,subTopicBodies,subTopicImages):
    powerpointName = '{}\{} Powerpoint.pptx'.format(topic.replace(' ','_'),topic)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0]) #title layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
        
    title.text = str(topic)
    if len(subTopicTitles) > 3:
        subtitleText = '{}, {}, {} and more'.format(*subTopicTitles[:3])
    else:
        subtitleText = ', '.join(subTopicTitles)
    #adding first few subtopics as title page subtitle

    subtitle.text = subtitleText

    leftText = True
    for topicTitle, body, images in zip(subTopicTitles, subTopicBodies, subTopicImages):
        leftText,presentation = addSlide(leftText,presentation,topicTitle,body,images)

    presentation = addRefsSlide(powerpointName,presentation,referencesContent)

    return presentation,powerpointName

def addSlide(leftText,presentation,topicTitle,body,images):
    print("ADDING SLIDE:", topicTitle)

    pictureSlide = bool(images)
    imageCount = len(images)

    if pictureSlide:
        bullet_slide_layout = presentation.slide_layouts[3] #picture and text layout
    else:
        bullet_slide_layout = presentation.slide_layouts[1] #just text layout

    slide = presentation.slides.add_slide(bullet_slide_layout)

    #add speaker notes below
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = body[:5000]

    #print(len(body))

    if len(body)>7500:
        start_chunk = body[:5000]
        end_chunk = body[-2500:]
        body = f"{start_chunk}\n\n[...CONTENT OMITTED...]\n\n{end_chunk}"

    #generate bullet points
    slideContent = makeBulletPoints(body)
    #print(len(slideContent))
    #print(len(slideContent.split('\n')))
    #print(len(body))

    lines = slideContent.split('\n')
    linesWithoutDashes = [line.lstrip('- ') for line in lines]
    slideContent = '\n'.join(linesWithoutDashes)

    if pictureSlide:
        leftText = not(leftText)
        fontSize = 18 #smaller font if image on slide
        if leftText:
            subtitle = slide.placeholders[1] # Text is left
            slide.shapes._spTree.remove(slide.placeholders[2]._element) # Remove right placeholder
        else:
            subtitle = slide.placeholders[2] # Text is right
            slide.shapes._spTree.remove(slide.placeholders[1]._element) # Remove left placeholder
    else:
        subtitle = slide.placeholders[1]
        fontSize = 24 #larger text if no images

    if len(slideContent) > 700:
        fontSize-=4
    elif len(slideContent) > 550:
        fontSize-=2
    elif len(slideContent) < 400:
        fontSize+=2

    if len(slideContent.split('\n'))<=5:
        fontSize+=2
    elif len(slideContent.split('\n'))>=8:
        fontSize-=1

    title = slide.shapes.title
    title.text = topicTitle

    # Assign bullet points to the text box
    subtitle.text = slideContent
    text_frame = subtitle.text_frame
    for paragraph in text_frame.paragraphs:
        font = paragraph.font
        font.size = Pt(fontSize)

    if pictureSlide:
        # Define the image area boundaries (based on slide layout 3)
        if leftText:
            IMAGE_AREA_LEFT = Inches(5)
        else:
            IMAGE_AREA_LEFT = Inches(0.5)
        IMAGE_AREA_TOP = Inches(1.5)
        IMAGE_AREA_WIDTH = Inches(4.5)
        IMAGE_AREA_HEIGHT = Inches(5.5)
        
        # Calculate image placement based on 1 or 2 images
        
        if imageCount == 1:
            localImagePath1, captionText1 = images[0]
            
            # Use original getImageSize function, passing area limits
            width, height, widthLost, heightLost = getImageSize(localImagePath1, maxWidth=4.5, maxHeight=5.5)
            
            # Place picture centered in the allocated area
            slide.shapes.add_picture(localImagePath1, IMAGE_AREA_LEFT + Inches(widthLost/2), IMAGE_AREA_TOP + Inches(heightLost/2), Inches(width), Inches(height))
            
            if captionText1:
                caption_top = IMAGE_AREA_TOP + Inches(height) + Inches(heightLost/2) + Inches(0.1) # 0.1 inch gap below image
                caption_box = slide.shapes.add_textbox(IMAGE_AREA_LEFT, caption_top, IMAGE_AREA_WIDTH, Inches(0.5))
                text_frame = caption_box.text_frame
                text_frame.text = captionText1
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                text_frame.paragraphs[0].font.size = Pt(10)
                text_frame.paragraphs[0].font.italic = True
                
        elif imageCount == 2:
            # New Max height for each stacked image (half height minus small margin)
            MAX_H_STACKED = 2.6
            GAP = 0.2
            
            localImagePath1, captionText1 = images[0]
            localImagePath2, captionText2 = images[1]
            
            # Get sizes constrained by the half-height area
            width1, height1, widthLost1, heightLost1 = getImageSize(localImagePath1, maxWidth=4.5, maxHeight=MAX_H_STACKED)
            width2, height2, widthLost2, heightLost2 = getImageSize(localImagePath2, maxWidth=4.5, maxHeight=MAX_H_STACKED)
            
            # Place Image 1 (Top)
            # Vertically center in the top half (1.5" to 1.5"+MAX_H_STACKED)
            top_y1 = IMAGE_AREA_TOP + Inches(heightLost1/2)
            slide.shapes.add_picture(localImagePath1, IMAGE_AREA_LEFT + Inches(widthLost1/2), top_y1, Inches(width1), Inches(height1))
            
            if captionText1:
                caption_top1 = top_y1 + Inches(height1) + Inches(0.1)
                caption_box1 = slide.shapes.add_textbox(IMAGE_AREA_LEFT, caption_top1-Inches(0.1), IMAGE_AREA_WIDTH, Inches(0.3))
                text_frame1 = caption_box1.text_frame
                text_frame1.text = captionText1
                text_frame1.paragraphs[0].alignment = PP_ALIGN.CENTER
                text_frame1.paragraphs[0].font.size = Pt(8)
                text_frame1.paragraphs[0].font.italic = True

            # Place Image 2 (Bottom)
            # Starts right below the top half area (1.5" + MAX_H_STACKED + GAP)
            top_y2 = IMAGE_AREA_TOP + Inches(MAX_H_STACKED) + Inches(GAP) + Inches(heightLost2/2)
            slide.shapes.add_picture(localImagePath2, IMAGE_AREA_LEFT + Inches(widthLost2/2), top_y2+Inches(0.1), Inches(width2), Inches(height2))

            if captionText2:
                caption_top2 = top_y2 + Inches(height2) + Inches(0.1)
                caption_box2 = slide.shapes.add_textbox(IMAGE_AREA_LEFT, caption_top2, IMAGE_AREA_WIDTH, Inches(0.3))
                text_frame2 = caption_box2.text_frame
                text_frame2.text = captionText2
                text_frame2.paragraphs[0].alignment = PP_ALIGN.CENTER
                text_frame2.paragraphs[0].font.size = Pt(8)
                text_frame2.paragraphs[0].font.italic = True

    return leftText,presentation

def addRefsSlide(powerpointName,presentation,referencesContent):
    print("ADDING SLIDE: References")
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = 'References'

    body = slide.placeholders[1]  # main body placeholder
    tf = body.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.level = 0
    run1 = p.add_run()
    run1.text = "Main source: "
    run2 = p.add_run()
    run2.text = url
    run2.hyperlink.address = url
    run1.font.size = Pt(20)
    run2.font.size = Pt(20)

    p = tf.add_paragraph()
    p.level = 0
    p.text = referencesContent
    p.font.size = Pt(12)

    os.makedirs(os.path.dirname(powerpointName), exist_ok=True)
    presentation.save(powerpointName)


#-------------------------------------------------------------------------------------------------------------------

ai = OpenAI(api_key=os.environ["OPENAI_API_KEY"])
headers = {"User-Agent": "Mozilla/5.0"}

subTopicTitles = []
subTopicBodies = []
subTopicImages = []

avoidedContents = ['Citations','Notes','See also','Sources','Further reading','External links','Gallery','Bibliography','Works cited','Collaborators','References','References 2']
avoidedImages = [
    'Question_book-new.svg',
    'Nuvola_apps_kaboodle.svg',
    'Ambox_current_red_Asia_Australia.svg',
    'Information_icon4.svg',
    'Climate_change_icon',
    'Symbol_list_class.svg',
    'Ambox_rewrite.svg',
    'Ambox_important.svg',
    'Wiki_letter_w_cropped.svg',
    'Commons-logo.svg',
    'Wikibooks-logo-en-noslogan.svg',
    'Semi-protection-shackle-keyhole.svg',
    'Wiki_letter_w.svg',
    'Red_flag_II.svg',
    'A_coloured_voting_box.svg',
    'Symbol-hammer-and-sickle.svg',
    '40px']

if __name__ == "__main__":
    url,topic,parentFolder = searchWikipedia()
    response = requests.get(url, headers=headers)
    html = response.text

    print('Generating Powerpoint!')
    start=time.time()
    referencesContent,subTopicTitles,subTopicBodies,subTopicImages = splitContent(html)
    presentation,powerpointName = generatePP(topic,referencesContent,subTopicTitles,subTopicBodies,subTopicImages)
    end = time.time()
    print(f'PowerPoint Generated! ({round(end-start, 1)} seconds)')

    os.startfile(powerpointName)

